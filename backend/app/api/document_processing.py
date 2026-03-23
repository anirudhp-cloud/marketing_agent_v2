"""
POST /api/brand/upload — Single file handling:
  1. Document parsing (PDF / DOCX / PPTX / image)
  2. Logo identification via GPT-4o vision
  3. Brand insights extraction via GPT-4o
"""
from __future__ import annotations

import base64
import io
import json
import logging
from pathlib import Path

from fastapi import APIRouter, UploadFile, File, HTTPException
from PIL import Image

from app.config import get_settings
from app.utils.llm import get_openai_client

logger = logging.getLogger(__name__)

router = APIRouter()

UPLOAD_DIR = Path(__file__).resolve().parent.parent.parent / "static" / "uploads"

ALLOWED_EXTENSIONS = {"pdf", "docx", "pptx", "png", "jpg", "jpeg", "webp"}

# ─── GPT-4o prompts ──────────────────────────────────────────────────────────

LOGO_SYSTEM_PROMPT = """You are a brand-design expert specialising in logo identification.
You must distinguish between:
- A LOGO: a company/brand mark — typically a symbol, wordmark, icon, emblem, or monogram used for brand identity. Logos are usually simple, vector-like, and appear on a solid or transparent background.
- NOT a logo: product photos, lifestyle images, illustrations, decorative graphics, screenshots, full-page renders, or photographs of any kind.

Analyse the image and determine:
1. Is this a company/brand logo? (yes/no) — be strict. If it's a photo, product image, or complex illustration, say no.
2. If yes, what type? (wordmark / icon / emblem / combination mark / monogram / abstract)
3. If yes, describe the logo briefly (shape, colours, text on it).
4. Confidence score 0-100 (be conservative — only score >70 if clearly a logo).
5. Does the logo have a transparent or solid-colour background?

Respond ONLY in JSON:
{"is_logo": true/false, "logo_type": "...", "description": "...", "confidence": 0-100, "has_simple_background": true/false}
"""

INSIGHTS_SYSTEM_PROMPT = """You are a senior brand strategist. Given the text extracted from a brand guidelines document, extract structured brand information.

Respond ONLY in JSON with these keys (use empty string or empty array if not found):
{
  "company_name": "string",
  "tagline": "string",
  "industry": "string",
  "product_category": "string",
  "description": "string",
  "target_audience": "string",
  "brand_personality": "string",
  "brand_values": ["string"],
  "brand_colours": ["#hex"],
  "typography": "string",
  "logo_description": "string",
  "competitors": ["string"],
  "unique_selling_points": ["string"],
  "tone_of_voice": "string",
  "key_messages": ["string"],
  "social_media_presence": "string"
}

Rules:
- For brand_colours, extract ALL hex colour codes mentioned. Look for primary, secondary, accent colours.
- For typography, note font families and their usage (headings, body, etc).
- Be thorough — extract every useful detail.
"""

# ─── Document extraction ─────────────────────────────────────────────────────


def _extract_from_pdf(file_bytes: bytes) -> tuple[str, list[tuple[bytes, str]]]:
    import fitz
    import re

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    text_parts: list[str] = []
    images: list[tuple[bytes, str]] = []

    for page in doc:
        # --- Structured text extraction (preserves table layout) ---
        page_text = page.get_text("text")

        # If plain text is empty/very short, try extracting as layout-preserving dict
        if len(page_text.strip()) < 20:
            # Render page to image and run OCR via PyMuPDF's built-in Tesseract bridge
            try:
                ocr_text = page.get_text("text", flags=fitz.TEXT_DEHYPHENATE)
                if len(ocr_text.strip()) < 20:
                    # Fallback: render page to pixmap for OCR-like extraction
                    pix = page.get_pixmap(dpi=200)
                    img_bytes_page = pix.tobytes("png")
                    images.append((img_bytes_page, "png"))
                    # Try extracting text from the rendered image via Tesseract if available
                    try:
                        import pytesseract
                        pil_img = Image.open(io.BytesIO(img_bytes_page))
                        ocr_result = pytesseract.image_to_string(pil_img)
                        if ocr_result.strip():
                            page_text = ocr_result
                    except ImportError:
                        logger.debug("pytesseract not installed — skipping OCR for scanned page")
                    except Exception as ocr_err:
                        logger.warning("OCR failed for page: %s", ocr_err)
                else:
                    page_text = ocr_text
            except Exception:
                pass

        # --- Table extraction via PyMuPDF ---
        try:
            tables = page.find_tables()
            for table in tables:
                extracted = table.extract()
                if extracted:
                    # Format table rows as pipe-separated for readable text
                    table_lines = []
                    for row in extracted:
                        cells = [str(c).strip() if c else "" for c in row]
                        table_lines.append(" | ".join(cells))
                    page_text += "\n\n[TABLE]\n" + "\n".join(table_lines) + "\n[/TABLE]\n"
        except Exception as table_err:
            logger.debug("Table extraction failed on page: %s", table_err)

        text_parts.append(page_text)

        # --- Extract embedded raster images ---
        for img_info in page.get_images(full=True):
            xref = img_info[0]
            try:
                base_image = doc.extract_image(xref)
                if base_image and base_image["image"]:
                    images.append((base_image["image"], base_image.get("ext", "png")))
            except Exception:
                continue

    # --- Extract hex colour codes from entire text ---
    full_text = "\n".join(text_parts)
    hex_colors = re.findall(r"#(?:[0-9a-fA-F]{3}){1,2}\b", full_text)
    if hex_colors:
        full_text += "\n\n[COLOUR CODES FOUND]: " + ", ".join(sorted(set(hex_colors)))

    doc.close()
    return full_text, images


def _extract_from_docx(file_bytes: bytes) -> tuple[str, list[tuple[bytes, str]]]:
    from docx import Document
    from docx.oxml.ns import qn
    import re

    doc = Document(io.BytesIO(file_bytes))
    text_parts: list[str] = []
    images: list[tuple[bytes, str]] = []
    style_info: list[str] = []

    # --- Extract body paragraphs with style/formatting hints ---
    for para in doc.paragraphs:
        line = para.text.strip()
        if not line:
            continue
        # Capture font info from runs for brand typography extraction
        fonts_in_para = set()
        has_bold = False
        has_italic = False
        for run in para.runs:
            if run.font.name:
                fonts_in_para.add(run.font.name)
            if run.bold:
                has_bold = True
            if run.italic:
                has_italic = True
        # Tag headings and styled text for better GPT extraction
        style_name = para.style.name if para.style else ""
        if "Heading" in style_name:
            line = f"## {line}"
        if fonts_in_para:
            style_info.append(f"Font '{', '.join(fonts_in_para)}' used in: \"{line[:80]}\"")
        text_parts.append(line)

    # --- Extract ALL tables (critical for brand guidelines) ---
    for table in doc.tables:
        table_lines: list[str] = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace("\n", " ")
                cells.append(cell_text)
            table_lines.append(" | ".join(cells))
        if table_lines:
            text_parts.append("\n[TABLE]\n" + "\n".join(table_lines) + "\n[/TABLE]")

    # --- Extract headers and footers ---
    for section in doc.sections:
        for header in (section.header, section.first_page_header):
            if header and header.paragraphs:
                header_text = " ".join(p.text.strip() for p in header.paragraphs if p.text.strip())
                if header_text:
                    text_parts.append(f"[HEADER] {header_text}")
        for footer in (section.footer, section.first_page_footer):
            if footer and footer.paragraphs:
                footer_text = " ".join(p.text.strip() for p in footer.paragraphs if p.text.strip())
                if footer_text:
                    text_parts.append(f"[FOOTER] {footer_text}")

    # --- Extract text boxes / floating text frames ---
    try:
        body = doc.element.body
        # Text boxes live inside w:txbxContent elements
        for txbx in body.iter(qn("w:txbxContent")):
            for p_elem in txbx.iter(qn("w:t")):
                if p_elem.text and p_elem.text.strip():
                    text_parts.append(p_elem.text.strip())
        # Also check for fallback content in mc:AlternateContent
        for fallback in body.iter(qn("mc:Fallback")):
            for p_elem in fallback.iter(qn("w:t")):
                if p_elem.text and p_elem.text.strip():
                    text_parts.append(p_elem.text.strip())
    except Exception as e:
        logger.debug("Text-box extraction encountered issue: %s", e)

    # --- Extract theme colours from document XML ---
    try:
        theme_part = doc.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        if theme_part:
            theme_xml = theme_part._element.xml if hasattr(theme_part, '_element') else ""
            # Extract colour values from theme XML (a:srgbClr val="RRGGBB")
            srgb_colors = re.findall(r'val="([0-9A-Fa-f]{6})"', theme_xml)
            if srgb_colors:
                unique_colors = sorted(set(f"#{c}" for c in srgb_colors))
                text_parts.append(f"\n[THEME COLOURS]: {', '.join(unique_colors)}")
    except Exception:
        pass

    # --- Append typography summary ---
    if style_info:
        text_parts.append("\n[TYPOGRAPHY]\n" + "\n".join(sorted(set(style_info))) + "\n[/TYPOGRAPHY]")

    # --- Extract hex colour codes from text ---
    full_text = "\n".join(text_parts)
    hex_colors = re.findall(r"#(?:[0-9a-fA-F]{3}){1,2}\b", full_text)
    if hex_colors:
        full_text += "\n\n[COLOUR CODES FOUND]: " + ", ".join(sorted(set(hex_colors)))

    # --- Extract images from relationships ---
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                blob = rel.target_part.blob
                ext = rel.target_part.content_type.split("/")[-1].replace("jpeg", "jpg")
                images.append((blob, ext))
            except Exception:
                continue

    return full_text, images


def _extract_from_pptx(file_bytes: bytes) -> tuple[str, list[tuple[bytes, str]]]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    text_parts: list[str] = []
    images: list[tuple[bytes, str]] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_parts.append(shape.text_frame.text)
            if shape.shape_type == 13:  # Picture
                ext = shape.image.content_type.split("/")[-1].replace("jpeg", "jpg")
                images.append((shape.image.blob, ext))

    return "\n".join(text_parts), images


_EXTRACTORS = {
    "pdf": _extract_from_pdf,
    "docx": _extract_from_docx,
    "pptx": _extract_from_pptx,
}


def _filter_viable_images(
    images: list[tuple[bytes, str]],
    min_size: int = 50,
    max_images: int = 10,
) -> list[tuple[bytes, str]]:
    """Drop tiny images (icons/bullets), score by logo-likelihood, and cap."""
    scored: list[tuple[float, bytes, str]] = []
    for img_bytes, ext in images:
        try:
            img = Image.open(io.BytesIO(img_bytes))
            w, h = img.width, img.height
            if w < min_size or h < min_size:
                continue

            # Score images by logo-likelihood:
            # - Logos tend to be moderate size (not huge full-page renders)
            # - Logos tend to have aspect ratios near 1:1 to ~3:1
            # - Logos often have limited unique colours
            # - Very large images (>2000px) are likely photos, not logos
            score = 50.0
            aspect = max(w, h) / max(min(w, h), 1)

            # Prefer moderate-sized images (100-800px range)
            if 80 <= max(w, h) <= 800:
                score += 25
            elif max(w, h) > 2000:
                score -= 20  # likely a full-page photo

            # Prefer reasonable aspect ratios (logos are rarely ultra-wide/tall)
            if aspect <= 4.0:
                score += 15
            elif aspect > 8.0:
                score -= 15  # probably a banner decoration

            # Check if image has transparency (alpha channel = likely a logo)
            if img.mode in ("RGBA", "LA", "PA"):
                score += 20

            # Check colour simplicity — logos tend to have fewer unique colours
            try:
                small = img.copy()
                small.thumbnail((64, 64))
                colours = small.convert("RGB").getcolors(maxcolors=512)
                if colours and len(colours) < 64:
                    score += 15  # simple palette = more likely a logo
            except Exception:
                pass

            scored.append((score, img_bytes, ext))
        except Exception:
            continue

    # Sort by logo-likelihood score (descending), not raw file size
    scored.sort(key=lambda x: x[0], reverse=True)
    return [(img_bytes, ext) for _, img_bytes, ext in scored[:max_images]]


# ─── GPT-4o vision: logo identification ──────────────────────────────────────


async def _identify_logos(
    images: list[tuple[bytes, str]],
) -> list[tuple[int, str, int]]:
    """
    Send each image to GPT-4o vision to check if it's a logo.
    Returns [(image_index, description, confidence)] sorted by confidence desc.
    Uses asyncio.gather to check all images in parallel.
    """
    if not images:
        return []

    client = get_openai_client()
    settings = get_settings()

    async def _check_one(idx: int, img_bytes: bytes, ext: str) -> tuple[int, str, int] | None:
        mime = f"image/{'jpeg' if ext in ('jpg', 'jpeg') else ext}"
        b64 = base64.b64encode(img_bytes).decode("ascii")
        try:
            resp = await client.chat.completions.create(
                model=settings.azure_openai_deployment,
                messages=[
                    {"role": "system", "content": LOGO_SYSTEM_PROMPT},
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": (
                                    "Is this image a company/brand logo? "
                                    "Remember: product photos, lifestyle images, and illustrations are NOT logos. "
                                    "Only brand marks, wordmarks, icons, emblems, and monograms count."
                                ),
                            },
                            {
                                "type": "image_url",
                                "image_url": {"url": f"data:{mime};base64,{b64}", "detail": "high"},
                            },
                        ],
                    },
                ],
                max_tokens=300,
                temperature=0,
            )
            raw = (resp.choices[0].message.content or "").strip()
            if raw.startswith("```"):
                raw = raw.split("\n", 1)[-1].rsplit("```", 1)[0]
            parsed = json.loads(raw)
            if parsed.get("is_logo") and parsed.get("confidence", 0) >= 50:
                return (idx, parsed.get("description", ""), parsed.get("confidence", 50))
        except Exception as e:
            logger.warning("Logo check failed for image %d: %s", idx, e)
        return None

    import asyncio
    tasks = [_check_one(idx, img_bytes, ext) for idx, (img_bytes, ext) in enumerate(images)]
    results_raw = await asyncio.gather(*tasks)
    results = [r for r in results_raw if r is not None]
    results.sort(key=lambda x: x[2], reverse=True)
    return results


# ─── GPT-4o: brand insights extraction ───────────────────────────────────────


async def _extract_brand_insights(text: str) -> dict:
    """Send document text to GPT-4o to extract structured brand insights."""
    if not text.strip():
        return {}

    client = get_openai_client()
    settings = get_settings()
    truncated = text[:24000]

    try:
        resp = await client.chat.completions.create(
            model=settings.azure_openai_deployment,
            messages=[
                {"role": "system", "content": INSIGHTS_SYSTEM_PROMPT},
                {"role": "user", "content": f"Here is the brand document text:\n\n{truncated}"},
            ],
            max_tokens=1500,
            temperature=0,
        )
        raw = (resp.choices[0].message.content or "").strip()
        if raw.startswith("```"):
            raw = raw.split("\n", 1)[-1].rsplit("```", 1)[0]
        return json.loads(raw)
    except Exception as e:
        logger.error("Brand insights extraction failed: %s", e)
        return {}


# ─── Logo post-processing ─────────────────────────────────────────────────────


def _postprocess_logo(img: Image.Image) -> Image.Image:
    """Clean up an extracted logo: trim whitespace, remove solid background, ensure RGBA.

    If background removal makes the image mostly transparent (>90% of pixels),
    we fall back to the original image with only whitespace trimmed.
    """
    img = img.convert("RGBA")
    original = img.copy()

    # Try to remove solid-colour background by checking corner pixels
    pixels = img.load()
    w, h = img.size
    corners = [
        pixels[0, 0],
        pixels[w - 1, 0],
        pixels[0, h - 1],
        pixels[w - 1, h - 1],
    ]

    # If 3+ corners share the same colour, treat it as background
    from collections import Counter
    corner_counts = Counter(corners)
    bg_color, count = corner_counts.most_common(1)[0]
    if count >= 3:
        # Make matching pixels transparent (with tolerance)
        tolerance = 30
        transparent_count = 0
        total_pixels = w * h
        for y in range(h):
            for x in range(w):
                r, g, b, a = pixels[x, y]
                if (
                    abs(r - bg_color[0]) < tolerance
                    and abs(g - bg_color[1]) < tolerance
                    and abs(b - bg_color[2]) < tolerance
                ):
                    pixels[x, y] = (r, g, b, 0)
                    transparent_count += 1

        # If >90% of pixels were made transparent, the logo content was lost
        if total_pixels > 0 and (transparent_count / total_pixels) > 0.90:
            logger.warning(
                "Background removal made %.0f%% of pixels transparent — "
                "falling back to original image (bg_color=%s)",
                (transparent_count / total_pixels) * 100,
                bg_color,
            )
            img = original

    # Trim transparent border
    bbox = img.getbbox()
    if bbox:
        img = img.crop(bbox)
    else:
        logger.warning("Logo image is fully transparent after processing — using original")
        img = original
        bbox = img.getbbox()
        if bbox:
            img = img.crop(bbox)

    return img


# ─── Route ────────────────────────────────────────────────────────────────────


def _safe_name(filename: str) -> str:
    stem = filename.rsplit(".", 1)[0]
    safe = "".join(c if c.isalnum() or c in "-_" else "_" for c in stem)
    return safe[:60]


@router.post("/upload")
async def upload_brand_document(file: UploadFile = File(...)):
    """
    Upload a brand guidelines file → parse document → identify logo via
    GPT-4o vision → extract brand insights → return results.
    """
    if not file.filename:
        raise HTTPException(400, "No file provided")

    ext = file.filename.rsplit(".", 1)[-1].lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(400, f"Unsupported file type: .{ext}")

    file_bytes = await file.read()
    if not file_bytes:
        raise HTTPException(400, "Empty file")

    # 0. Save the original uploaded file to disk
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    original_filename = f"{_safe_name(file.filename)}.{ext}"
    (UPLOAD_DIR / original_filename).write_bytes(file_bytes)
    logger.info("Saved original file: %s (%d bytes)", original_filename, len(file_bytes))

    # 1. Extract text + images from document
    try:
        if ext in _EXTRACTORS:
            text, raw_images = _EXTRACTORS[ext](file_bytes)
        elif ext in ("png", "jpg", "jpeg", "webp"):
            text, raw_images = "", [(file_bytes, ext)]
        else:
            raise HTTPException(400, f"Unsupported: .{ext}")
    except HTTPException:
        raise
    except Exception as parse_err:
        logger.error("Failed to parse %s file: %s", ext, parse_err)
        raise HTTPException(
            422,
            f"Could not parse the .{ext} file. It may be corrupted, "
            f"password-protected, or in an unsupported format. Error: {parse_err}"
        )

    images = _filter_viable_images(raw_images)
    logger.info("Extracted %d images (%d viable), %d chars text",
                len(raw_images), len(images), len(text))

    # 2. Run logo identification and brand insights in parallel
    import asyncio
    extracted_filenames: list[str] = []
    fallback_name: str | None = None

    if images:
        # Save largest image as fallback
        img_bytes_0, img_ext_0 = images[0]
        fallback_name = f"{_safe_name(file.filename)}_image.{img_ext_0}"
        (UPLOAD_DIR / fallback_name).write_bytes(img_bytes_0)
        logger.info("Saved largest image as fallback: %s", fallback_name)

    # Run both GPT-4o calls concurrently
    logos, insights = await asyncio.gather(
        _identify_logos(images) if images else asyncio.sleep(0, result=[]),
        _extract_brand_insights(text),
        return_exceptions=True,
    )

    # Handle logo results
    if isinstance(logos, Exception):
        logger.error("Logo identification failed: %s", logos)
        logos = []

    if images and logos:
        best_idx, logo_desc, confidence = logos[0]
        img_bytes, img_ext = images[best_idx]

        # Post-process logo: trim whitespace, ensure transparency, save as PNG
        try:
            logo_img = Image.open(io.BytesIO(img_bytes))
            logo_img = _postprocess_logo(logo_img)
            buf = io.BytesIO()
            logo_img.save(buf, format="PNG")
            processed_bytes = buf.getvalue()
            logo_filename = f"{_safe_name(file.filename)}_logo.png"
            (UPLOAD_DIR / logo_filename).write_bytes(processed_bytes)
            extracted_filenames.append(logo_filename)
            logger.info("Logo saved: %s (confidence=%d, desc=%s)",
                        logo_filename, confidence, logo_desc)
        except Exception as e:
            logger.warning("Logo post-processing failed, saving raw: %s", e)
            logo_filename = f"{_safe_name(file.filename)}_logo.{img_ext}"
            (UPLOAD_DIR / logo_filename).write_bytes(img_bytes)
            extracted_filenames.append(logo_filename)
    else:
        logger.info("No logo identified in uploaded document")

    # Handle insights results
    if isinstance(insights, Exception):
        logger.error("Brand insights extraction failed: %s", insights)
        insights = {}

    return {
        "insights": insights,
        "extracted_images": extracted_filenames,
        "original_file": original_filename,
    }
